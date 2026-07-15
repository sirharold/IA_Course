from pathlib import Path

from pptx import Presentation
from pptx.util import Emu


ROOT = Path(__file__).resolve().parent
SOURCE_PPTX = ROOT.parent / "clase 6" / "Clase_6_desde_clase_5.pptx"
AGENDA_IMAGE = ROOT / "imagenes" / "01-agenda-clase-07.png"
SLIDE_IMAGES_DIR = ROOT / "imagenes"
OUTPUT_PPTX = ROOT / "Clase_7_desde_clase_6.pptx"


def delete_slides_after(prs: Presentation, keep_count: int) -> None:
    slide_ids = list(prs.slides._sldIdLst)
    for slide_id in slide_ids[keep_count:]:
        rel_id = slide_id.rId
        prs.part.drop_rel(rel_id)
        prs.slides._sldIdLst.remove(slide_id)


def update_title_slide(prs: Presentation) -> None:
    slide = prs.slides[0]
    replacements = {
        "GITHUB COPILOT CLI Y AUTOMATIZACIÓN AVANZADA": "SKILLS, CUSTOMIZACION REUSABLE Y MCP EN COPILOT",
        "CLASE 6 DE 9 · CAPACITACIÓN TÉCNICA": "CLASE 7 DE 9 · CAPACITACIÓN TÉCNICA",
    }

    for shape in slide.shapes:
        if not hasattr(shape, "text_frame") or not shape.text.strip():
            continue
        text = shape.text.strip()
        if text in replacements:
            shape.text_frame.paragraphs[0].runs[0].text = replacements[text]


def replace_agenda_slide_with_image(prs: Presentation) -> None:
    slide = prs.slides[1]
    slide.shapes.add_picture(
        str(AGENDA_IMAGE),
        Emu(0),
        Emu(0),
        width=prs.slide_width,
        height=prs.slide_height,
    )


def append_image_slides(prs: Presentation) -> None:
    blank_layout = prs.slide_layouts[6]
    for image_path in sorted(SLIDE_IMAGES_DIR.glob("*.png")):
        if image_path.name == AGENDA_IMAGE.name:
            continue
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(image_path),
            Emu(0),
            Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )


def main() -> None:
    prs = Presentation(str(SOURCE_PPTX))
    delete_slides_after(prs, keep_count=3)
    update_title_slide(prs)
    replace_agenda_slide_with_image(prs)
    append_image_slides(prs)
    prs.save(str(OUTPUT_PPTX))
    print(OUTPUT_PPTX)


if __name__ == "__main__":
    main()
